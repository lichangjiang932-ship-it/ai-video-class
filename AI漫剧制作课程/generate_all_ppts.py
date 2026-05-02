#!/usr/bin/env python3
"""
AI漫剧制作课程 PPT v8 — 视觉升级版
渐变背景 + 多彩配色 + 装饰几何 + 动画入场 + 精致排版
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree

# ══════════════════════════════════════════════
# 配色方案 — 多彩渐变
# ══════════════════════════════════════════════

# 深色背景系
BG_DARK   = RGBColor(0x0A, 0x0A, 0x14)
BG_CARD   = RGBColor(0x14, 0x16, 0x22)
BG_CARD2  = RGBColor(0x1A, 0x1D, 0x2C)
BG_HOVER  = RGBColor(0x22, 0x25, 0x38)

# 文字色
T1 = RGBColor(0xF0, 0xF0, 0xF5)   # 主标题 - 纯白
T2 = RGBColor(0xB8, 0xBA, 0xC8)   # 正文 - 柔灰
T3 = RGBColor(0x6E, 0x70, 0x82)   # 辅助 - 暗灰
T4 = RGBColor(0x44, 0x46, 0x58)   # 装饰 - 极暗

# 分割线
LINE  = RGBColor(0x2A, 0x2D, 0x3E)
LINE2 = RGBColor(0x35, 0x38, 0x4A)

# ── 模块强调色（每个模块独立配色，更丰富）──
MODULE_PALETTES = [
    # 0: 课程导论 - 星空蓝渐变
    {'primary': RGBColor(0x38, 0x7A, 0xFF), 'secondary': RGBColor(0x6C, 0x5C, 0xFF),
     'accent': RGBColor(0x00, 0xD4, 0xFF), 'glow': RGBColor(0x1A, 0x3A, 0x6E)},
    # 1: 行业认知 - 琥珀金
    {'primary': RGBColor(0xFF, 0x9F, 0x1A), 'secondary': RGBColor(0xFF, 0x6B, 0x35),
     'accent': RGBColor(0xFF, 0xD9, 0x3D), 'glow': RGBColor(0x4A, 0x2A, 0x0A)},
    # 2: 剧本创作 - 翡翠绿
    {'primary': RGBColor(0x00, 0xC9, 0x7B), 'secondary': RGBColor(0x00, 0xE6, 0x96),
     'accent': RGBColor(0x38, 0xFF, 0xBD), 'glow': RGBColor(0x0A, 0x3A, 0x2A)},
    # 3: 分镜设计 - 活力橙
    {'primary': RGBColor(0xFF, 0x6B, 0x35), 'secondary': RGBColor(0xFF, 0x8C, 0x42),
     'accent': RGBColor(0xFF, 0xB3, 0x47), 'glow': RGBColor(0x4A, 0x1A, 0x0A)},
    # 4: 即梦AI - 梦幻紫
    {'primary': RGBColor(0x8A, 0x4E, 0xF0), 'secondary': RGBColor(0xA8, 0x55, 0xF7),
     'accent': RGBColor(0xC0, 0x84, 0xFF), 'glow': RGBColor(0x2A, 0x1A, 0x4A)},
    # 5: 其他图像 - 珊瑚粉
    {'primary': RGBColor(0xFF, 0x5C, 0x8A), 'secondary': RGBColor(0xFF, 0x7B, 0xA3),
     'accent': RGBColor(0xFF, 0xA0, 0xBC), 'glow': RGBColor(0x4A, 0x1A, 0x2A)},
    # 6: 可灵AI - 电光蓝
    {'primary': RGBColor(0x00, 0x96, 0xFF), 'secondary': RGBColor(0x38, 0xB6, 0xFF),
     'accent': RGBColor(0x64, 0xD2, 0xFF), 'glow': RGBColor(0x0A, 0x2A, 0x4A)},
    # 7: Seedance - 薰衣草紫
    {'primary': RGBColor(0x7C, 0x3A, 0xED), 'secondary': RGBColor(0x9F, 0x67, 0xFF),
     'accent': RGBColor(0xC4, 0xB5, 0xFD), 'glow': RGBColor(0x2A, 0x1A, 0x4A)},
    # 8: 海螺AI - 青碧色
    {'primary': RGBColor(0x06, 0xB6, 0xD4), 'secondary': RGBColor(0x22, 0xD3, 0xEE),
     'accent': RGBColor(0x67, 0xE8, 0xF9), 'glow': RGBColor(0x0A, 0x3A, 0x3A)},
    # 9: 角色一致性 - 玫瑰金
    {'primary': RGBColor(0xF4, 0x3F, 0x5E), 'secondary': RGBColor(0xFB, 0x71, 0x85),
     'accent': RGBColor(0xFD, 0xA4, 0xAF), 'glow': RGBColor(0x4A, 0x0A, 0x1A)},
    # 10: 音频制作 - 金橙色
    {'primary': RGBColor(0xF5, 0x9E, 0x0B), 'secondary': RGBColor(0xFA, 0xBB, 0x31),
     'accent': RGBColor(0xFD, 0xE0, 0x47), 'glow': RGBColor(0x4A, 0x3A, 0x0A)},
    # 11: 剪辑后期 - 森林绿
    {'primary': RGBColor(0x10, 0xB9, 0x81), 'secondary': RGBColor(0x34, 0xD3, 0x99),
     'accent': RGBColor(0x6E, 0xE7, 0xB7), 'glow': RGBColor(0x0A, 0x3A, 0x2A)},
    # 12: 一站式平台 - 天际蓝
    {'primary': RGBColor(0x38, 0x7A, 0xFF), 'secondary': RGBColor(0x60, 0xA5, 0xFA),
     'accent': RGBColor(0x93, 0xC5, 0xFD), 'glow': RGBColor(0x0A, 0x2A, 0x4A)},
    # 13: ComfyUI - 科技紫
    {'primary': RGBColor(0xA8, 0x55, 0xF7),
     'secondary': RGBColor(0xC0, 0x84, 0xFF), 'accent': RGBColor(0xDD, 0xD6, 0xFE),
     'glow': RGBColor(0x2A, 0x1A, 0x4A)},
    # 14: 工业化 - 钢铁灰蓝
    {'primary': RGBColor(0x64, 0x74, 0x8B), 'secondary': RGBColor(0x94, 0xA3, 0xB8),
     'accent': RGBColor(0xCB, 0xD5, 0xE1), 'glow': RGBColor(0x1A, 0x2A, 0x3A)},
    # 15: 分发变现 - 烈焰红
    {'primary': RGBColor(0xEF, 0x44, 0x44), 'secondary': RGBColor(0xF8, 0x71, 0x71),
     'accent': RGBColor(0xFC, 0xA5, 0xA5), 'glow': RGBColor(0x4A, 0x0A, 0x0A)},
    # 16: 合规规范 - 深海蓝
    {'primary': RGBColor(0x1E, 0x40, 0xAF), 'secondary': RGBColor(0x38, 0x5E, 0xE0),
     'accent': RGBColor(0x60, 0xA5, 0xFA), 'glow': RGBColor(0x0A, 0x1A, 0x3A)},
    # 17: 综合实战 - 彩虹渐变(用金色代表)
    {'primary': RGBColor(0xFF, 0xB8, 0x00), 'secondary': RGBColor(0xFF, 0x8C, 0x00),
     'accent': RGBColor(0xFF, 0xD7, 0x00), 'glow': RGBColor(0x4A, 0x3A, 0x0A)},
]

SW, SH = Inches(13.333), Inches(7.5)
OUT = '/root/.openclaw/workspace/ai-video-class/AI漫剧制作课程/PPT课件'


# ══════════════════════════════════════════════
# 工具函数
# ══════════════════════════════════════════════

def _pal(n):
    """获取模块配色"""
    if 0 <= n < len(MODULE_PALETTES):
        return MODULE_PALETTES[n]
    return MODULE_PALETTES[0]

def _dim(c, factor=0.35):
    return RGBColor(max(0, int(c[0]*factor)), max(0, int(c[1]*factor)), max(0, int(c[2]*factor)))

def _bright(c, add=60):
    return RGBColor(min(255, c[0]+add), min(255, c[1]+add), min(255, c[2]+add))

def _mix(c1, c2, ratio=0.5):
    """混合两种颜色"""
    return RGBColor(
        int(c1[0]*ratio + c2[0]*(1-ratio)),
        int(c1[1]*ratio + c2[1]*(1-ratio)),
        int(c1[2]*ratio + c2[2]*(1-ratio)),
    )


# ── 基础绘制 ──

def _bg(s, color=BG_DARK):
    f = s.background.fill; f.solid(); f.fore_color.rgb = color

def _rect(s, x, y, w, h, c, alpha=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    if alpha is not None:
        _set_alpha(sh, alpha)
    return sh

def _rrect(s, x, y, w, h, c, border=None, border_w=1):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(border_w)
    return sh

def _circ(s, x, y, sz, c, alpha=None):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    if alpha is not None:
        _set_alpha(sh, alpha)
    return sh

def _set_alpha(shape, alpha_pct):
    """设置形状透明度 (0-100, 100=完全透明)"""
    try:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            spPr = sp.find(qn('a:spPr'))
        if spPr is None:
            return
        solid = spPr.find(qn('a:solidFill'))
        if solid is not None:
            color_elem = solid[0]
            alpha_elem = etree.SubElement(color_elem, qn('a:alpha'))
            alpha_elem.set('val', str(int((100 - alpha_pct) * 1000)))
    except Exception:
        pass  # 透明度是装饰性功能，失败不影响主流程

def _t(s, x, y, w, h, txt, sz, c=T1, bold=False, al=PP_ALIGN.LEFT, name='Microsoft YaHei'):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = str(txt)
    p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = name; p.alignment = al
    return tb

def _t_multi(s, x, y, w, h, lines, default_sz=14, default_c=T2):
    """多行文字，每行可自定义格式"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, dict):
            p.text = line.get('text', '')
            p.font.size = Pt(line.get('sz', default_sz))
            p.font.color.rgb = line.get('c', default_c)
            p.font.bold = line.get('bold', False)
            p.font.name = line.get('font', 'Microsoft YaHei')
            p.alignment = line.get('al', PP_ALIGN.LEFT)
            p.space_after = Pt(line.get('space_after', 4))
        else:
            p.text = str(line)
            p.font.size = Pt(default_sz); p.font.color.rgb = default_c
            p.font.name = 'Microsoft YaHei'; p.space_after = Pt(4)
    return tb

def _line(s, x, y, w, c=LINE, th=Pt(1.5)):
    _rect(s, x, y, w, th, c)


# ── 高级装饰组件 ──

def _gradient_bar(s, x, y, w, h, c1, c2, steps=8):
    """水平渐变条"""
    sw = w // steps
    for i in range(steps):
        ratio = i / (steps - 1) if steps > 1 else 0
        c = _mix(c1, c2, 1 - ratio)
        _rect(s, x + i * sw, y, sw + (1 if i == 0 else 0), h, c)

def _glow_circle(s, x, y, sz, c, alpha=85):
    """发光圆：半透明大圆营造光晕"""
    _circ(s, x, y, sz, c, alpha)

def _dot_pattern(s, x, y, cols, rows, spacing, c, alpha=80):
    """点阵装饰"""
    for r in range(rows):
        for col_i in range(cols):
            _circ(s, x + col_i * spacing, y + r * spacing, Pt(3), c, alpha)

def _corner_deco(s, ac):
    """右上角几何装饰"""
    _rect(s, Inches(12.2), Inches(0), Inches(1.133), Pt(3), ac['primary'])
    _rect(s, Inches(13.1), Inches(0), Pt(3), Inches(0.8), ac['primary'])
    _circ(s, Inches(12.6), Inches(0.3), Inches(0.15), ac['accent'], 70)

def _bg_deco(s, ac):
    """背景装饰：大光晕 + 点阵"""
    _glow_circle(s, Inches(9), Inches(-1), Inches(5), ac['primary'], 92)
    _glow_circle(s, Inches(-2), Inches(4), Inches(4), ac['secondary'], 94)
    _dot_pattern(s, Inches(10.5), Inches(5.5), 6, 4, Inches(0.35), ac['accent'], 88)


# ── 头部/底部 ──

def _head(s, title, ac, icon=""):
    pal = ac if isinstance(ac, dict) else _pal(0)
    # 顶部渐变条
    _gradient_bar(s, Inches(0), Inches(0), SW, Pt(3), pal['primary'], pal['secondary'])
    # 标题背景
    _rect(s, Inches(0), Pt(3), SW, Inches(1.05), BG_CARD)
    # 左侧色条
    _rect(s, Inches(0), Pt(3), Pt(5), Inches(1.05), pal['primary'])
    # 标题
    label = f"{icon}  {title}" if icon else title
    _t(s, Inches(0.6), Inches(0.16), Inches(11), Inches(0.65), label, 24, T1, bold=True)
    # 底部分割线
    _line(s, Inches(0), Inches(1.18), SW, LINE, Pt(1))
    # 右侧小装饰
    _circ(s, Inches(12.5), Inches(0.25), Inches(0.12), pal['accent'], 60)

def _foot(s, mn, mname, pn, ac):
    pal = ac if isinstance(ac, dict) else _pal(0)
    _line(s, Inches(0), Inches(7.05), SW, LINE, Pt(1))
    _t(s, Inches(0.8), Inches(7.1), Inches(5), Inches(0.3), f"模块{mn:02d} · {mname}", 9, T3)
    _t(s, Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.3), str(pn), 9, T3, al=PP_ALIGN.RIGHT)
    # 页脚渐变小条
    _gradient_bar(s, Inches(0.8), Inches(7.05), Inches(1.5), Pt(2), pal['primary'], pal['secondary'])


# ── 卡片组件 ──

def _card(s, x, y, w, h, ac=None, glow=False):
    """精致卡片：圆角 + 可选顶部色条 + 可选光晕"""
    _rrect(s, x, y, w, h, BG_CARD, LINE, 1)
    if ac:
        pal = ac if isinstance(ac, dict) else _pal(0)
        _gradient_bar(s, x, y, w, Pt(3), pal['primary'], pal['secondary'])
    if glow and ac:
        pal = ac if isinstance(ac, dict) else _pal(0)
        _glow_circle(s, x + w//2 - Inches(1), y - Inches(0.5), Inches(2), pal['primary'], 93)

def _icon_badge(s, x, y, icon, ac, sz=28):
    """图标徽章：圆形背景 + emoji"""
    pal = ac if isinstance(ac, dict) else _pal(0)
    _circ(s, x, y, Inches(0.55), pal['glow'])
    _circ(s, x + Pt(3), y + Pt(3), Inches(0.48), _dim(pal['primary'], 0.3))
    _t(s, x + Pt(5), y + Pt(3), Inches(0.45), Inches(0.45), icon, sz, T1, al=PP_ALIGN.CENTER)

def _numbered_item(s, x, y, num, text, ac, sz=14):
    pal = ac if isinstance(ac, dict) else _pal(0)
    # 编号圆
    _circ(s, x, y + Pt(2), Inches(0.3), pal['primary'])
    _t(s, x + Pt(1), y + Pt(2), Inches(0.28), Inches(0.28), f"{num:02d}", 10, T1, bold=True, al=PP_ALIGN.CENTER)
    _t(s, x + Inches(0.4), y, Inches(11), Inches(0.5), text, sz, T2)

def _bullet_item(s, x, y, text, ac, sz=14):
    pal = ac if isinstance(ac, dict) else _pal(0)
    # 渐变色块
    _gradient_bar(s, x, y + Inches(0.08), Inches(0.06), Inches(0.4), pal['primary'], pal['secondary'])
    _t(s, x + Inches(0.22), y, Inches(11.5), Inches(0.6), text, sz, T2)

def _tag(s, x, y, text, ac, sz=10):
    """标签"""
    pal = ac if isinstance(ac, dict) else _pal(0)
    w = Inches(len(text) * 0.12 + 0.3)
    _rrect(s, x, y, w, Inches(0.3), pal['glow'])
    _t(s, x, y, w, Inches(0.3), text, sz, pal['accent'], bold=True, al=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 动画支持
# ══════════════════════════════════════════════

def _add_fade_in(shape, delay_ms=0, duration_ms=500):
    """为形状添加淡入动画"""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return

    # 创建动画元素
    timing = None
    # 查找或创建 timing 元素
    slide = sp.getroottree().getroot()
    timing_list = slide.findall('.//' + qn('p:timing'))
    if timing_list:
        timing = timing_list[0]
    else:
        timing = etree.SubElement(slide.find(qn('p:cSld')).getparent(), qn('p:timing'))

    # 简化：直接在 slide XML 中添加 entrance animation
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = etree.SubElement(timing, qn('p:tnLst'))

    # 创建 par (parallel) 容器
    par = etree.SubElement(tnLst, qn('p:par'))
    cTn = etree.SubElement(par, qn('p:cTn'))
    cTn.set('id', '1')
    cTn.set('dur', 'indefinite')
    cTn.set('restart', 'never')
    cTn.set('nodeType', 'tmRoot')

    childTnLst = etree.SubElement(cTn, qn('p:childTnLst'))
    seq = etree.SubElement(childTnLst, qn('p:seq'))
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    seqCtn = etree.SubElement(seq, qn('p:cTn'))
    seqCtn.set('id', '2')
    seqCtn.set('dur', 'indefinite')
    seqCtn.set('nodeType', 'mainSeq')

    seqChildTnLst = etree.SubElement(seqCtn, qn('p:childTnLst'))

    # par for this animation
    anim_par = etree.SubElement(seqChildTnLst, qn('p:par'))
    anim_cTn = etree.SubElement(anim_par, qn('p:cTn'))
    anim_cTn.set('id', '3')
    anim_cTn.set('fill', 'hold')

    # set delay
    stCondLst = etree.SubElement(anim_cTn, qn('p:stCondLst'))
    stCond = etree.SubElement(stCondLst, qn('p:cond'))
    stCond.set('delay', str(delay_ms))

    anim_childTnLst = etree.SubElement(anim_cTn, qn('p:childTnLst'))

    # set (effect)
    set_elem = etree.SubElement(anim_childTnLst, qn('p:set'))
    set_cBhvr = etree.SubElement(set_elem, qn('p:cBhvr'))

    set_cTn = etree.SubElement(set_cBhvr, qn('p:cTn'))
    set_cTn.set('id', '4')
    set_cTn.set('dur', '1')
    set_cTn.set('fill', 'hold')

    set_stCondLst = etree.SubElement(set_cTn, qn('p:stCondLst'))
    set_stCond = etree.SubElement(set_stCondLst, qn('p:cond'))
    set_stCond.set('delay', '0')

    # attrNameLst
    attrNameLst = etree.SubElement(set_cBhvr, qn('p:attrNameLst'))
    attrName = etree.SubElement(attrNameLst, qn('p:attrName'))
    attrName.text = 'style.visibility'

    # to value
    to = etree.SubElement(set_elem, qn('p:to'))
    strVal = etree.SubElement(to, qn('p:strVal'))
    strVal.set('val', 'visible')

    # set spid
    spTgt = etree.SubElement(set_cBhvr, qn('p:spTgt'))
    spTgt.set('spid', str(sp.get('id', '')))

    # animateEffect (fade)
    effect_elem = etree.SubElement(anim_childTnLst, qn('p:effect'))
    effect_elem.set('transition', 'in')
    effect_elem.set('filter', 'fade')

    effect_cBhvr = etree.SubElement(effect_elem, qn('p:cBhvr'))
    effect_cTn = etree.SubElement(effect_cBhvr, qn('p:cTn'))
    effect_cTn.set('id', '5')
    effect_cTn.set('dur', str(duration_ms))

    effect_stCondLst = etree.SubElement(effect_cTn, qn('p:stCondLst'))
    effect_stCond = etree.SubElement(effect_stCondLst, qn('p:cond'))
    effect_stCond.set('delay', '0')

    effect_spTgt = etree.SubElement(effect_cBhvr, qn('p:spTgt'))
    effect_spTgt.set('spid', str(sp.get('id', '')))

    effect_attrNameLst = etree.SubElement(effect_cBhvr, qn('p:attrNameLst'))
    effect_attrName = etree.SubElement(effect_attrNameLst, qn('p:attrName'))
    effect_attrName.text = 'style.visibility'

    # prevCondLst for seq
    prevCondLst = etree.SubElement(seq, qn('p:prevCondLst'))
    prevCond = etree.SubElement(prevCondLst, qn('p:cond'))
    prevCond.set('evt', 'onPrev')
    prevCond.set('delay', '0')
    tgtEl = etree.SubElement(prevCond, qn('p:tgtEl'))
    sldTgt = etree.SubElement(tgtEl, qn('p:sldTgt'))

    nextCondLst = etree.SubElement(seq, qn('p:nextCondLst'))
    nextCond = etree.SubElement(nextCondLst, qn('p:cond'))
    nextCond.set('evt', 'onNext')
    nextCond.set('delay', '0')
    tgtEl2 = etree.SubElement(nextCond, qn('p:tgtEl'))
    sldTgt2 = etree.SubElement(tgtEl2, qn('p:sldTgt'))


def _add_slide_transition(s, transition_type='fade', duration_ms=700, advance_ms=None):
    """为幻灯片添加切换动画"""
    transition = etree.SubElement(s._element, qn('p:transition'))
    transition.set('spd', 'med')
    transition.set('advClick', '1')

    if advance_ms:
        transition.set('advTm', str(advance_ms))

    if transition_type == 'fade':
        fade = etree.SubElement(transition, qn('p:fade'))
    elif transition_type == 'push':
        push = etree.SubElement(transition, qn('p:push'))
        push.set('dir', 'l')
    elif transition_type == 'wipe':
        wipe = etree.SubElement(transition, qn('p:wipe'))
    elif transition_type == 'split':
        split = etree.SubElement(transition, qn('p:split'))
        split.set('orient', 'horz')
        split.set('dir', 'out')


# ══════════════════════════════════════════════
# 幻灯片类型
# ══════════════════════════════════════════════

def make_cover(prs, title, subtitle, mn=0):
    """封面：大渐变背景 + 几何装饰 + 光晕"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)

    # 背景装饰
    _glow_circle(s, Inches(8), Inches(-2), Inches(8), ac['primary'], 93)
    _glow_circle(s, Inches(-3), Inches(3), Inches(6), ac['secondary'], 95)
    _glow_circle(s, Inches(5), Inches(5), Inches(4), ac['accent'], 96)

    # 左侧大色块
    _rect(s, Inches(0), Inches(0), Inches(5.8), SH, ac['primary'])
    _rect(s, Inches(5.8), Inches(0), Inches(0.2), SH, ac['secondary'])
    _rect(s, Inches(6.0), Inches(0), Inches(0.1), SH, _dim(ac['secondary'], 0.5))

    # 左侧装饰：点阵
    _dot_pattern(s, Inches(0.5), Inches(5.5), 8, 3, Inches(0.3), T1, 80)
    # 左侧大编号
    if mn > 0:
        _t(s, Inches(0.3), Inches(0.8), Inches(5), Inches(3), f"{mn:02d}", 140, RGBColor(0xFF,0xFF,0xFF), bold=True)
        # 模块标签
        _rrect(s, Inches(0.5), Inches(4.0), Inches(2.0), Inches(0.36), BG_DARK)
        _t(s, Inches(0.5), Inches(4.0), Inches(2.0), Inches(0.36),
           "M O D U L E", 10, ac['accent'], al=PP_ALIGN.CENTER)
        _gradient_bar(s, Inches(0.5), Inches(4.5), Inches(2.5), Pt(3), ac['accent'], ac['primary'])
    else:
        _t(s, Inches(1.0), Inches(1.5), Inches(3.5), Inches(2.5), "🎬", 100, T1, al=PP_ALIGN.CENTER)
        _t(s, Inches(0.5), Inches(4.0), Inches(5), Inches(0.5), "AI漫剧制作", 20, RGBColor(0xFF,0xFF,0xFF), bold=True, al=PP_ALIGN.CENTER)

    # 右侧内容
    # 课程标签
    _rrect(s, Inches(6.6), Inches(1.2), Inches(6.0), Inches(0.5), ac['glow'])
    _t(s, Inches(6.6), Inches(1.2), Inches(6.0), Inches(0.5),
       "AI漫剧制作全流程课程 · 2026版", 12, ac['accent'], al=PP_ALIGN.CENTER)

    # 标题
    _t(s, Inches(6.6), Inches(2.2), Inches(6.0), Inches(2.0), title, 38, T1, bold=True)

    # 渐变装饰线
    _gradient_bar(s, Inches(6.6), Inches(4.2), Inches(3.5), Pt(3), ac['primary'], ac['secondary'])

    # 副标题
    _t(s, Inches(6.6), Inches(4.6), Inches(6.0), Inches(0.8), subtitle, 15, T2)

    # 底部信息卡片
    _card(s, Inches(6.6), Inches(5.8), Inches(6.0), Inches(0.7))
    for i, st in enumerate(["124课时", "17模块", "8-16周", "全流程"]):
        x = Inches(6.75) + Inches(i * 1.5)
        _t(s, x, Inches(5.8), Inches(1.3), Inches(0.7), st, 11, ac['primary'], bold=True, al=PP_ALIGN.CENTER)
        if i < 3:
            _rect(s, x + Inches(1.3), Inches(6.0), Pt(1), Inches(0.3), LINE)

    # 右上角装饰
    _corner_deco(s, ac)

    # 添加淡入动画
    _add_slide_transition(s, 'fade', 500)


def make_toc(prs, items, mn, mname, pn):
    """目录：编号卡片 + 渐变装饰"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _bg_deco(s, ac)
    _head(s, "本节内容", ac, "📋")

    col1 = items[:len(items)//2 + len(items)%2]
    col2 = items[len(items)//2 + len(items)%2:]

    for i, item in enumerate(col1):
        y = Inches(1.5) + i * Inches(0.62)
        if y > Inches(6.4): break
        _numbered_item(s, Inches(0.8), y, i + 1, item, ac)

    for i, item in enumerate(col2):
        y = Inches(1.5) + i * Inches(0.62)
        if y > Inches(6.4): break
        _numbered_item(s, Inches(7.0), y, len(col1) + i + 1, item, ac)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_content(prs, title, bullets, mn, mname, pn):
    """内容页：渐变色块标记 + 呼吸间距"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _bg_deco(s, ac)
    _head(s, title, ac)

    n = len(bullets)
    for i, b in enumerate(bullets):
        y = Inches(1.5) + i * Inches(0.88)
        if y > Inches(6.4): break
        _bullet_item(s, Inches(0.8), y, b, ac, 15)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_stat_cards(prs, title, stats, mn, mname, pn):
    """数据页：大数字 + 渐变装饰 + 光晕"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _head(s, title, ac, "📊")

    n = min(len(stats), 4)
    cw = Inches(12.3 / n - 0.2)
    gap = Inches(0.25)
    sx = Inches(0.5)

    for i, st in enumerate(stats[:n]):
        x = sx + i * (cw + gap)
        # 卡片 + 光晕
        _card(s, x, Inches(1.5), cw, Inches(5.2), ac, glow=True)
        # 渐变顶部条
        _gradient_bar(s, x, Inches(1.5), cw, Pt(4), ac['primary'], ac['secondary'])

        # 大数字
        _t(s, x + Inches(0.1), Inches(1.9), cw - Inches(0.2), Inches(1.5),
           st.get('value', ''), 52, ac['primary'], bold=True, al=PP_ALIGN.CENTER)

        # 单位
        if st.get('unit'):
            _t(s, x + Inches(0.1), Inches(3.2), cw - Inches(0.2), Inches(0.3),
               st['unit'], 13, T3, al=PP_ALIGN.CENTER)

        # 分割线
        _line(s, x + Inches(0.3), Inches(3.65), cw - Inches(0.6), LINE, Pt(1))

        # 标签
        _t(s, x + Inches(0.1), Inches(3.85), cw - Inches(0.2), Inches(0.4),
           st.get('label', ''), 15, T1, bold=True, al=PP_ALIGN.CENTER)

        # 描述
        if st.get('desc'):
            _t(s, x + Inches(0.15), Inches(4.4), cw - Inches(0.3), Inches(0.8),
               st['desc'], 11, T3, al=PP_ALIGN.CENTER)

        # 趋势标签
        if st.get('trend'):
            is_up = any(c in st['trend'] for c in ['↑', '+'])
            tc = RGBColor(0x00, 0xCE, 0x7E) if is_up else RGBColor(0xFF, 0x2E, 0x5E)
            _rrect(s, x + Inches(0.2), Inches(5.5), cw - Inches(0.4), Inches(0.35), BG_CARD2)
            _t(s, x + Inches(0.2), Inches(5.5), cw - Inches(0.4), Inches(0.35),
               st['trend'], 12, tc, bold=True, al=PP_ALIGN.CENTER)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_two_col(prs, title, lt, li, rt, ri, mn, mname, pn):
    """双栏对比：渐变卡片 + 光晕"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _bg_deco(s, ac)
    _head(s, title, ac)

    # 左栏
    _card(s, Inches(0.5), Inches(1.5), Inches(5.85), Inches(5.2), ac)
    _t(s, Inches(0.8), Inches(1.75), Inches(5.2), Inches(0.45), lt, 18, ac['primary'], bold=True)
    _gradient_bar(s, Inches(0.8), Inches(2.3), Inches(1.8), Pt(2), ac['primary'], ac['secondary'])
    for i, item in enumerate(li):
        y = Inches(2.55) + i * Inches(0.58)
        if y > Inches(6.3): break
        _circ(s, Inches(0.8), y + Inches(0.06), Inches(0.16), ac['primary'])
        _t(s, Inches(1.1), y, Inches(5.0), Inches(0.45), item, 13, T2)

    # 右栏（用 secondary 色）
    _card(s, Inches(6.85), Inches(1.5), Inches(5.85), Inches(5.2))
    _gradient_bar(s, Inches(6.85), Inches(1.5), Inches(5.85), Pt(3), ac['secondary'], ac['accent'])
    _t(s, Inches(7.15), Inches(1.75), Inches(5.2), Inches(0.45), rt, 18, ac['secondary'], bold=True)
    _gradient_bar(s, Inches(7.15), Inches(2.3), Inches(1.8), Pt(2), ac['secondary'], ac['accent'])
    for i, item in enumerate(ri):
        y = Inches(2.55) + i * Inches(0.58)
        if y > Inches(6.3): break
        _circ(s, Inches(7.15), y + Inches(0.06), Inches(0.16), ac['secondary'])
        _t(s, Inches(7.45), y, Inches(5.0), Inches(0.45), item, 13, T2)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_table(prs, title, headers, rows, mn, mname, pn):
    """表格：渐变表头 + 交替行色"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _bg_deco(s, ac)
    _head(s, title, ac)

    cols = len(headers); nr = len(rows) + 1
    tw = Inches(12.1); rh = Inches(0.55)
    ts = s.shapes.add_table(nr, cols, Inches(0.6), Inches(1.4), tw, rh * nr)
    tbl = ts.table; cw = tw / cols
    for j in range(cols): tbl.columns[j].width = int(cw)

    # 表头：渐变色
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = ac['primary']
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(12); p.font.color.rgb = T1; p.font.bold = True
            p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
        c.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(val); c.fill.solid()
            c.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else BG_CARD2
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(11); p.font.color.rgb = T2
                p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
            c.vertical_anchor = MSO_ANCHOR.MIDDLE

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_key_point(prs, title, kp, exp, mn, mname, pn):
    """重点页：大字居中 + 光晕背景 + 渐变装饰"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)

    # 大光晕
    _glow_circle(s, Inches(3), Inches(0), Inches(7), ac['primary'], 93)
    _glow_circle(s, Inches(0), Inches(3), Inches(5), ac['secondary'], 95)

    # 渐变装饰线
    _gradient_bar(s, Inches(3.5), Inches(1.2), Inches(6.3), Pt(3), ac['primary'], ac['secondary'])

    # 核心大字
    _t(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.5),
       kp, 30, T1, bold=True, al=PP_ALIGN.CENTER)

    # 中间装饰
    _gradient_bar(s, Inches(5.5), Inches(4.3), Inches(2.3), Pt(2), ac['primary'], ac['secondary'])

    # 说明
    _t(s, Inches(1.5), Inches(4.7), Inches(10.3), Inches(1.8),
       exp, 15, T2, al=PP_ALIGN.CENTER)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 500)


def make_steps(prs, title, steps, mn, mname, pn):
    """步骤页：卡片式步骤 + 连接箭头"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _bg_deco(s, ac)
    _head(s, title, ac)

    n = len(steps)
    if n <= 4:
        cw = Inches(12.3 / n - 0.2)
        gap = Inches(0.25)
        for i, step in enumerate(steps):
            x = Inches(0.5) + i * (cw + gap)
            _card(s, x, Inches(1.5), cw, Inches(5.0), ac)
            # 大编号
            _t(s, x + Inches(0.2), Inches(1.8), cw - Inches(0.4), Inches(0.9),
               f"{i+1}", 48, ac['primary'], bold=True, al=PP_ALIGN.CENTER)
            _gradient_bar(s, x + Inches(0.3), Inches(2.8), cw - Inches(0.6), Pt(2), ac['primary'], ac['secondary'])
            _t(s, x + Inches(0.2), Inches(3.0), cw - Inches(0.4), Inches(3.3), step, 13, T2)
            # 箭头
            if i < n - 1:
                _t(s, x + cw + Inches(0.01), Inches(2.8), Inches(0.2), Inches(0.3),
                   "›", 18, ac['primary'], bold=True, al=PP_ALIGN.CENTER)
    else:
        for i, step in enumerate(steps):
            y = Inches(1.5) + i * Inches(0.85)
            if y > Inches(6.4): break
            _bullet_item(s, Inches(0.8), y, step, ac, 14)
            _t(s, Inches(0.8), y, Inches(0.3), Inches(0.5), f"{i+1}", 9, T1, bold=True)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_practice(prs, title, tasks, mn, mname, pn):
    """练习页：checkbox样式"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _bg_deco(s, ac)

    # 用绿色系
    green = {'primary': RGBColor(0x00, 0xCE, 0x7E), 'secondary': RGBColor(0x34, 0xD3, 0x99),
             'accent': RGBColor(0x6E, 0xE7, 0xB7), 'glow': RGBColor(0x0A, 0x3A, 0x2A)}
    _head(s, title, green, "🛠️")

    for i, task in enumerate(tasks):
        y = Inches(1.5) + i * Inches(0.85)
        if y > Inches(6.4): break
        # checkbox
        _rrect(s, Inches(0.8), y + Inches(0.08), Inches(0.32), Inches(0.32), BG_CARD, green['primary'])
        _gradient_bar(s, Inches(0.8), y + Inches(0.08), Inches(0.32), Pt(2), green['primary'], green['secondary'])
        _t(s, Inches(1.3), y, Inches(0.5), Inches(0.5), f"#{i+1}", 11, green['primary'], bold=True)
        _t(s, Inches(1.8), y, Inches(10.5), Inches(0.6), task, 14, T2)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_summary(prs, items, mn, mname, pn):
    """总结页"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _bg_deco(s, ac)
    _head(s, "本节小结", ac, "📝")

    for i, item in enumerate(items):
        y = Inches(1.5) + i * Inches(0.78)
        if y > Inches(6.4): break
        _numbered_item(s, Inches(0.8), y, i + 1, item, ac, 14)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_section(prs, stitle, sub, mn, mname, pn):
    """章节页：极简 + 大光晕 + 渐变"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)

    # 大光晕
    _glow_circle(s, Inches(2), Inches(0), Inches(8), ac['primary'], 92)
    _glow_circle(s, Inches(-2), Inches(3), Inches(5), ac['secondary'], 94)

    # 左侧色条
    _rect(s, Inches(0), Inches(0), Inches(0.08), SH, ac['primary'])

    # 渐变装饰线
    _gradient_bar(s, Inches(1.2), Inches(2.7), Inches(5), Pt(3), ac['primary'], ac['secondary'])

    _t(s, Inches(1.2), Inches(3.1), Inches(10), Inches(1.0), stitle, 34, T1, bold=True)
    _t(s, Inches(1.2), Inches(4.2), Inches(10), Inches(0.6), sub, 15, T3)

    # 右下装饰
    _dot_pattern(s, Inches(10), Inches(5.5), 5, 3, Inches(0.3), ac['accent'], 85)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 500)


def make_timeline(prs, title, events, mn, mname, pn):
    """时间线"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _bg_deco(s, ac)
    _head(s, title, ac, "⏱️")

    n = len(events)
    if n == 0:
        _foot(s, mn, mname, pn, ac); return

    ly = Inches(3.5)
    _gradient_bar(s, Inches(1.0), ly, Inches(11.3), Pt(3), ac['primary'], ac['secondary'])

    step = Inches(11.3) / n
    for i, ev in enumerate(events):
        x = Inches(1.0) + i * step
        _circ(s, x + Inches(0.35), ly - Inches(0.1), Inches(0.22), ac['primary'])
        _circ(s, x + Inches(0.38), ly - Inches(0.07), Inches(0.16), ac['accent'])
        label = ev.get('label', '') if isinstance(ev, dict) else str(ev)
        _t(s, x, ly - Inches(0.65), Inches(1.2), Inches(0.4), label, 14, ac['primary'], bold=True, al=PP_ALIGN.CENTER)
        desc = ev.get('desc', '') if isinstance(ev, dict) else ''
        if desc:
            dy = Inches(1.6) if i % 2 == 0 else Inches(4.0)
            _card(s, x - Inches(0.05), dy, Inches(1.5), Inches(1.2), ac)
            _t(s, x + Inches(0.05), dy + Inches(0.15), Inches(1.3), Inches(0.9), desc, 11, T2, al=PP_ALIGN.CENTER)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_icon_grid(prs, title, items, mn, mname, pn):
    """图标网格：卡片 + 渐变 + 光晕"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)
    _bg_deco(s, ac)
    _head(s, title, ac)

    cols = 3; cw = Inches(3.7); ch = Inches(2.0); gap = Inches(0.35)
    sx, sy = Inches(0.65), Inches(1.4)

    for i, item in enumerate(items[:9]):
        col, row = i % cols, i // cols
        x = sx + col * (cw + gap); y = sy + row * (ch + gap)
        _card(s, x, y, cw, ch, ac)
        icon = item.get('icon', '📌') if isinstance(item, dict) else '📌'
        _t(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5), icon, 24, ac['primary'])
        it = item.get('title', '') if isinstance(item, dict) else item
        _t(s, x + Inches(0.2), y + Inches(0.75), cw - Inches(0.4), Inches(0.35), it, 14, T1, bold=True)
        if isinstance(item, dict) and item.get('desc'):
            _t(s, x + Inches(0.2), y + Inches(1.2), cw - Inches(0.4), Inches(0.65), item['desc'], 11, T3)

    _foot(s, mn, mname, pn, ac)
    _add_slide_transition(s, 'fade', 400)


def make_end(prs, nt, mn, mname):
    """结束页：大光晕 + 渐变"""
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
    ac = _pal(mn)

    # 大光晕
    _glow_circle(s, Inches(3), Inches(1), Inches(7), ac['primary'], 92)

    _gradient_bar(s, Inches(4.5), Inches(1.5), Inches(4.3), Pt(3), ac['primary'], ac['secondary'])

    _t(s, Inches(0), Inches(2.2), SW, Inches(1.0), "✓", 64, ac['primary'], bold=True, al=PP_ALIGN.CENTER)
    _t(s, Inches(0), Inches(3.5), SW, Inches(0.7), "本模块结束", 30, T1, bold=True, al=PP_ALIGN.CENTER)
    _gradient_bar(s, Inches(5.5), Inches(4.4), Inches(2.3), Pt(2), ac['primary'], ac['secondary'])
    if nt:
        _t(s, Inches(0), Inches(4.7), SW, Inches(0.5), f"下一模块：{nt}", 16, T2, al=PP_ALIGN.CENTER)
    _t(s, Inches(0), Inches(6.5), SW, Inches(0.4), "AI漫剧制作全流程课程 · 2026版", 10, T3, al=PP_ALIGN.CENTER)

    _add_slide_transition(s, 'fade', 600)


def make_pr(prs):
    prs.slide_width = SW; prs.slide_height = SH; return prs


# ══════════════════════════════════════════════
# 模块定义（内置模块0和1）
# ══════════════════════════════════════════════

MODULES = [
    {
        'num': 0, 'name': '课程导论',
        'title': 'AI漫剧制作全流程课程',
        'subtitle': '124课时 · 17个模块 · 从零基础到独立制作完整AI漫剧',
        'next': '行业认知与趋势',
        'pages': [
            {'type': 'toc', 'items': [
                '行业认知与趋势', '大语言模型与剧本创作', '分镜设计与提示词工程',
                '即梦AI深度精通', '其他图像生成工具', '可灵AI深度精通',
                'Seedance 2.0深度精通', '海螺AI与Vidu', '角色一致性攻克',
                '配音配乐与音画同步', '剪辑与后期全流程', '一站式平台精讲',
                'ComfyUI工作流', '工业化生产与团队协作', '平台分发与商业变现',
                '合规要求与行业规范', '综合实战项目'
            ]},
            {'type': 'content', 'title': '课程目标', 'bullets': [
                '掌握AI漫剧从剧本到成片的完整制作流程',
                '精通即梦AI、可灵AI、Seedance 2.0等核心工具',
                '能独立制作一集完整的AI漫剧',
                '了解平台分发与商业变现的路径',
                '建立工业化生产思维，实现高效量产'
            ]},
            {'type': 'stat_cards', 'title': '行业数据一览', 'stats': [
                {'value': '240', 'unit': '亿元', 'label': '市场规模', 'trend': '↑ +50% YoY'},
                {'value': '2.8', 'unit': '亿', 'label': '用户规模', 'trend': '↑ +40% YoY'},
                {'value': '15', 'unit': '亿次/天', 'label': '日均播放', 'trend': '↑ +50% YoY'},
                {'value': '120', 'unit': '万', 'label': '创作者', 'trend': '↑ +50% YoY'},
            ]},
            {'type': 'two_col', 'title': '学习路径选择',
             'left_title': '🅰️ 零基础速成（8周）', 'left_items': ['快速上手出作品', '重点模块：1,2,3,4,6,10,11,12,17', '适合：想快速入行变现'],
             'right_title': '🅱️ 专业进阶（16周）', 'right_items': ['全面掌握全链路', '全部17个模块', '适合：追求专业深度']},
            {'type': 'step', 'title': '课程学习路线图', 'steps': [
                '行业认知 → 理解AI漫剧是什么、市场多大',
                '创作技能 → 剧本、分镜、提示词工程',
                '工具精通 → 即梦AI、可灵AI、Seedance等',
                '后期制作 → 音频、剪辑、调色、成片',
                '商业运营 → 平台分发、变现、合规',
                '综合实战 → 从0到1独立完成一集AI漫剧'
            ]},
            {'type': 'key_point', 'title': '核心竞争力',
             'key_point': 'AI漫剧的核心竞争力\n= 创意 × 工具熟练度 × 量产能力',
             'explanation': '工具会更新迭代，但创意能力和工业化思维是持久的竞争力。本课程不仅教工具操作，更注重培养你的创作思维和生产效率。'},
        ]
    },
    {
        'num': 1, 'name': '行业认知与趋势',
        'title': '模块一：AI漫剧行业认知与趋势',
        'subtitle': '6课时 · 从行业全景到创业路径',
        'next': '大语言模型与剧本创作',
        'pages': [
            {'type': 'toc', 'items': ['行业全景：从手工作坊到智能流水线', '市场数据：240亿规模与2.8亿用户', '内容形态：沙雕漫→动态漫→AI原生漫剧', '五步工业化流程', '爆款拆解：《斩仙台》等', '创业路径：个人vs团队vs工作室']},
            {'type': 'content', 'title': '什么是AI漫剧？', 'bullets': [
                '利用AI技术辅助或主导制作的动态漫画短剧',
                '制作周期：传统1-3个月/集 → AI 1-3天/集',
                '制作成本：传统5-20万/集 → AI 500-5000元/集',
                '团队规模：传统10-30人 → AI 1-5人',
                '产能：传统2-4集/月 → AI 30+集/月'
            ]},
            {'type': 'timeline', 'title': '行业发展四个阶段', 'events': [
                {'label': '2023', 'desc': '沙雕漫\n静态图+配音'},
                {'label': '2024', 'desc': '动态漫\n图生视频应用'},
                {'label': '2025', 'desc': 'AI原生\n全链路AI辅助'},
                {'label': '2026', 'desc': 'AI仿真人\n接近真人质量'},
            ]},
            {'type': 'stat_cards', 'title': '2026年核心市场数据', 'stats': [
                {'value': '240', 'unit': '亿元', 'label': '市场规模', 'trend': '↑ +50%', 'desc': '较2025年160亿'},
                {'value': '2.8', 'unit': '亿', 'label': '用户规模', 'trend': '↑ +40%', 'desc': '较2025年2.0亿'},
                {'value': '15', 'unit': '亿次', 'label': '日均播放', 'trend': '↑ +50%', 'desc': '较2025年10亿次'},
                {'value': '120', 'unit': '万', 'label': '创作者', 'trend': '↑ +50%', 'desc': '较2025年80万'},
            ]},
            {'type': 'step', 'title': '五步工业化流程', 'steps': [
                '剧本创作：用AI生成分集剧本',
                '分镜设计：将剧本转化为分镜表',
                '文生图：用即梦AI生成关键帧画面',
                '图生视频：用可灵AI/Seedance让画面动起来',
                '剪辑配音：用剪映完成后期制作'
            ]},
            {'type': 'content', 'title': '爆款案例拆解', 'bullets': [
                '《斩仙台》：玄幻题材，角色一致性标杆，单集播放破千万',
                '《气运三角洲》：都市逆袭，节奏把控精准',
                '《霍去病》：历史题材，画面质感接近电影级别',
                '共同特点：前3秒强钩子、角色一致性好、节奏紧凑'
            ]},
            {'type': 'two_col', 'title': '创业路径选择',
             'left_title': '👤 个人创作者', 'left_items': ['月入1-10万', '日产1集', '月成本~5000元', '适合：副业/自由职业'],
             'right_title': '👥 小团队（3-5人）', 'right_items': ['月入10-50万', '日产4集', '月成本~7万', '适合：创业团队']},
            {'type': 'practice', 'title': '本模块实操任务', 'tasks': [
                '关注3个AI漫剧行业公众号/博主，建立信息源',
                '在抖音搜索10个AI漫剧账号，分析其内容特点',
                '画出完整的AI漫剧产业链地图',
                '确定自己的定位：个人/团队/工作室',
                '选择一个感兴趣的题材方向'
            ]},
        ]
    },
]


def generate_module_ppt(md):
    """生成单个模块PPT"""
    prs = Presentation(); make_pr(prs)
    num, name = md['num'], md['name']
    make_cover(prs, md['title'], md['subtitle'], num)
    pn = 1
    for page in md['pages']:
        t = page['type']
        if t == 'toc': make_toc(prs, page['items'], num, name, pn)
        elif t == 'content': make_content(prs, page['title'], page['bullets'], num, name, pn)
        elif t == 'two_col': make_two_col(prs, page['title'], page['left_title'], page['left_items'], page['right_title'], page['right_items'], num, name, pn)
        elif t == 'table': make_table(prs, page['title'], page['headers'], page['rows'], num, name, pn)
        elif t == 'stat_cards': make_stat_cards(prs, page['title'], page['stats'], num, name, pn)
        elif t == 'icon_grid': make_icon_grid(prs, page['title'], page['items'], num, name, pn)
        elif t == 'timeline': make_timeline(prs, page['title'], page['events'], num, name, pn)
        elif t == 'key_point': make_key_point(prs, page.get('title','核心要点'), page['key_point'], page['explanation'], num, name, pn)
        elif t == 'step': make_steps(prs, page['title'], page['steps'], num, name, pn)
        elif t == 'practice': make_practice(prs, page['title'], page['tasks'], num, name, pn)
        elif t == 'section': make_section(prs, page['title'], page['subtitle'], num, name, pn)
        elif t == 'summary': make_summary(prs, page['items'], num, name, pn)
        pn += 1
    make_end(prs, md.get('next',''), num, name)
    fn = f"模块{num:02d}-{name}.pptx" if num > 0 else "00-课程导论.pptx"
    fp = os.path.join(OUT, fn); prs.save(fp)
    print(f"✅ {fn}（{len(prs.slides)}页）")
    return fp


if __name__ == '__main__':
    os.makedirs(OUT, exist_ok=True)
    for m in MODULES: generate_module_ppt(m)
    print(f"\n🎉 已生成 {len(MODULES)} 个PPT文件")
